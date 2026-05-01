from pptx import Presentation
from pptx.util import Inches, Pt
import re

INPUT = "Instant Retail Co/O2O Strategy/08-slide-content.md"
OUTPUT = "Finals/presentation_final.pptx"

with open(INPUT, 'r', encoding='utf-8') as f:
    raw = f.read()

# Split slides by horizontal rule marker (---)
blocks = [b.strip() for b in re.split(r"\n---\n", raw) if b.strip()]

prs = Presentation()
# Choose a title+content layout index (commonly 1)
layout_title_content = prs.slide_layouts[1]

for block in blocks:
    # Only process slide blocks that start with '## Slide'
    if not block.startswith('## Slide'):
        continue

    # Extract short label (after '## Slide N — label')
    m = re.search(r"## Slide \d+ — (.*)", block)
    short_label = m.group(1).strip() if m else ''

    # Extract Eyebrow, Title, Hero Title, Body Text bullets
    eyebrow = re.search(r"- \*\*Eyebrow\*\*: (.*)", block)
    eyebrow = eyebrow.group(1).strip() if eyebrow else ''

    title = re.search(r"- \*\*Title\*\*: (.*)", block)
    if title:
        title_text = title.group(1).strip()
    else:
        # Try hero title multiline
        hero = re.search(r"- \*\*Hero Title\*\*:\n((?:\s+- \".*\"\n)+)", block)
        if hero:
            lines = re.findall(r'\- \"(.*)\"', hero.group(1))
            title_text = ' - '.join(lines)
        else:
            title_text = short_label

    # Extract body bullets (both "Body Text (bullets):" and Point X lines)
    bullets = []
    # Body Text bullets
    m_body = re.search(r"- \*\*Body Text\*\* \(bullets\):\n((?:\s+- .*\n)+)", block)
    if m_body:
        for line in m_body.group(1).splitlines():
            line = line.strip()
            if line.startswith('- '):
                bullets.append(line[2:].strip())
    # Point lines like '- **Point 1**: text'
    for p in re.findall(r"- \*\*Point \d+\*\*: (.*)", block):
        bullets.append(p.strip())

    # Fallback: any top-level '-' lines that look like content
    if not bullets:
        for line in block.splitlines():
            line = line.strip()
            if line.startswith('- ') and ':' not in line[:40]:
                bullets.append(line[2:].strip())

    # Build slide
    slide = prs.slides.add_slide(layout_title_content)
    try:
        slide.shapes.title.text = title_text
    except Exception:
        pass

    # Add eyebrow as small subtitle if present
    if eyebrow:
        # try to set subtitle placeholder if exists
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                try:
                    shape.text = eyebrow
                except Exception:
                    pass
                break

    # Add bullets to content placeholder
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    if bullets:
        first = tf.paragraphs[0]
        first.text = bullets[0]
        first.font.size = Pt(14)
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(14)
    else:
        tf.text = ''

# Ensure Finals directory exists (Presentation will create file path)
prs.save(OUTPUT)
print(f"Saved PPTX to {OUTPUT}")
